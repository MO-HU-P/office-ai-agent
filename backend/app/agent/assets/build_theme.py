"""同梱デザインテンプレート(default_theme.pptx)の生成スクリプト。

python-pptx標準の白紙テンプレート(4:3・Office 2007配色)をベースに、
16:9化・アプリのUI(frontend/src/index.css)と揃えたGoogle風テーマ配色・
レイアウト装飾を加えたテンプレートを生成する。

デザインの原則: 配色・フォント・装飾はこのテンプレート(データ)に持たせ、
エージェントは中身を流し込むだけにする。テーマ色を差し替えているため、
表(既定スタイルはaccent1)・図形(既定塗りはaccent1)・グラフ(系列色はaccent1〜6)の
既定の見た目もこのテンプレートに追従する。

再生成: docker compose run --rm --no-deps -v ./backend:/srv -w /srv backend \
        python app/agent/assets/build_theme.py
"""
import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Cm, Emu

# アプリUI(index.css)と同じパレット
BLUE = "1A73E8"      # accent1: 図形・表ヘッダー・グラフ第1系列
GREEN = "188038"     # accent2
AMBER = "F9AB00"     # accent3
RED = "D93025"       # accent4
PURPLE = "9334E6"    # accent5
TEAL = "12B5CB"      # accent6
TEXT = "202124"      # 本文(dk1)
TEXT2 = "3C4043"     # 補助(dk2)
SURFACE = "F8F9FA"   # 薄い背景(lt2)
BLUE_TINT = "E8F0FE" # 装飾用の薄い青(UIの選択色と同系)

SLIDE_W, SLIDE_H = Cm(33.867), Cm(19.05)  # 16:9


def _replace_theme(prs: Presentation) -> None:
    """テーマXMLの配色とフォントを差し替える(表・図形・グラフの既定色がここから決まる)。"""
    theme_part = prs.slide_master.part.part_related_by(RT.THEME)
    xml = theme_part.blob.decode("utf-8")
    replacements = [
        # 本文色・背景色(Office 2007テーマの既定値からの置換。値が変わったら要確認)
        ('<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>',
         f'<a:dk1><a:srgbClr val="{TEXT}"/></a:dk1>'),
        ('<a:dk2><a:srgbClr val="1F497D"/></a:dk2>', f'<a:dk2><a:srgbClr val="{TEXT2}"/></a:dk2>'),
        ('<a:lt2><a:srgbClr val="EEECE1"/></a:lt2>', f'<a:lt2><a:srgbClr val="{SURFACE}"/></a:lt2>'),
        ('<a:accent1><a:srgbClr val="4F81BD"/></a:accent1>', f'<a:accent1><a:srgbClr val="{BLUE}"/></a:accent1>'),
        ('<a:accent2><a:srgbClr val="C0504D"/></a:accent2>', f'<a:accent2><a:srgbClr val="{GREEN}"/></a:accent2>'),
        ('<a:accent3><a:srgbClr val="9BBB59"/></a:accent3>', f'<a:accent3><a:srgbClr val="{AMBER}"/></a:accent3>'),
        ('<a:accent4><a:srgbClr val="8064A2"/></a:accent4>', f'<a:accent4><a:srgbClr val="{RED}"/></a:accent4>'),
        ('<a:accent5><a:srgbClr val="4BACC6"/></a:accent5>', f'<a:accent5><a:srgbClr val="{PURPLE}"/></a:accent5>'),
        ('<a:accent6><a:srgbClr val="F79646"/></a:accent6>', f'<a:accent6><a:srgbClr val="{TEAL}"/></a:accent6>'),
        # フォント: 和文は游ゴシック(無ければ各環境の既定にフォールバック)
        ('<a:latin typeface="Calibri"/>', '<a:latin typeface="Yu Gothic"/>'),
        ('<a:ea typeface=""/>', '<a:ea typeface="游ゴシック"/>'),
    ]
    for old, new in replacements:
        assert old in xml, f"テーマXMLに想定値が見つかりません(python-pptx更新の影響?): {old}"
        xml = xml.replace(old, new)
    theme_part._blob = xml.encode("utf-8")


def _place(ph, left, top, width, height) -> None:
    ph.left, ph.top, ph.width, ph.height = Cm(left), Cm(top), Cm(width), Cm(height)


def _relayout_for_16x9(prs: Presentation) -> None:
    """4:3向けのプレースホルダー配置を16:9(左右余白1.7cm)に合わせて敷き直す。"""
    master = prs.slide_master
    for ph in master.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            _place(ph, 1.7, 0.9, 30.5, 2.5)
        elif idx == 1:
            _place(ph, 1.7, 3.9, 30.5, 13.3)
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            if idx in (10, 11, 12):  # 日付・フッター・ページ番号(通常非表示)は下端右寄せに
                _place(ph, {10: 1.7, 11: 12.0, 12: 27.9}[idx], 18.0, {10: 6.0, 11: 10.0, 12: 4.2}[idx], 0.8)
        name = layout.name
        phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
        if name == "Title Slide":
            _place(phs[0], 2.5, 6.6, 28.9, 3.6)   # 中央タイトル
            _place(phs[1], 2.5, 10.7, 28.9, 2.6)  # サブタイトル
        elif name == "Section Header":
            _place(phs[1], 2.6, 6.6, 28.7, 1.6)   # 上の小さな導入文
            _place(phs[0], 2.6, 8.4, 28.7, 3.6)   # 章タイトル
        elif 0 in phs:  # Title and Content / Title Only など通常レイアウト
            _place(phs[0], 1.7, 0.9, 30.5, 2.5)
            if 1 in phs:
                _place(phs[1], 1.7, 3.9, 30.5, 13.3)


def _add_bar(scratch, layout, left, top, width, height, color: str):
    """レイアウトに装飾の帯(枠線なしの塗りつぶし矩形)を追加する。
    LayoutShapesにはadd_shapeが無いため、作業用スライド上で作った図形の
    XML要素をレイアウトのspTree(最背面=プレースホルダーの背面)へ移植する。"""
    sp = scratch.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    el = sp._element
    el.getparent().remove(el)
    # nvGrpSpPr / grpSpPr の直後(=最背面)に挿入する
    layout.shapes._spTree.insert(2, el)


def _delete_slide(prs, index: int) -> None:
    """スライドをパーツ(rels)ごと削除する(sldIdLstから外すだけだとパーツが残り、
    保存時にzipへ重複エントリが書かれてしまう)。"""
    sldId = prs.slides._sldIdLst[index]
    prs.part.drop_rel(sldId.rId)
    prs.slides._sldIdLst.remove(sldId)


def _decorate(prs: Presentation) -> None:
    """各レイアウトにアクセント装飾を加える(色はテーマではなく固定値。装飾はデータに焼き込む)。"""
    w = SLIDE_W / 360000  # cm
    scratch = prs.slides.add_slide(prs.slide_layouts[6])  # 図形生成用の作業スライド
    for layout in prs.slide_layouts:
        name = layout.name
        if name == "Title Slide":
            # 下部に薄青の帯+その上端に青のライン
            _add_bar(scratch, layout, 0, 16.7, w, 19.05 - 16.7, BLUE_TINT)
            _add_bar(scratch, layout, 0, 16.7, w, 0.12, BLUE)
        elif name == "Section Header":
            # 薄い背景+章タイトル左の縦の青バー
            layout.background.fill.solid()
            layout.background.fill.fore_color.rgb = RGBColor.from_string(SURFACE)
            _add_bar(scratch, layout, 1.9, 8.5, 0.35, 3.2, BLUE)
        elif name == "Blank":
            continue  # 自由レイアウトは無地のまま
        else:
            # 通常レイアウト: 上端に青のライン
            _add_bar(scratch, layout, 0, 0, w, 0.18, BLUE)
    _delete_slide(prs, -1)


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    _replace_theme(prs)
    _relayout_for_16x9(prs)
    _decorate(prs)
    assert len(prs.slides) == 0, "テンプレートにスライドが残っています"
    prs.save(str(out_path))
    print(f"生成しました: {out_path}")


if __name__ == "__main__":
    build(Path(__file__).parent / "default_theme.pptx")
