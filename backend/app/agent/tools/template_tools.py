"""テンプレート差し込みツール: {{キー}} プレースホルダーの一括置換"""
import re
from typing import Any

from docx import Document
from langchain_core.tools import tool
from openpyxl import load_workbook
from pptx import Presentation

from ...atomic import atomic_save
from ...config import resolve_workspace_path

# {{ キー }} 形式(前後の空白は許容)。キーに { } は使えない
PLACEHOLDER_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")


def _apply_text(text: str, values: dict[str, Any], counter: dict[str, int], missing: set[str]) -> str:
    def _sub(m: re.Match) -> str:
        key = m.group(1)
        if key in values:
            counter[key] = counter.get(key, 0) + 1
            return str(values[key])
        missing.add(key)
        return m.group(0)

    return PLACEHOLDER_RE.sub(_sub, text)


def _replace_in_paragraph(p, values: dict[str, Any], counter: dict[str, int], missing: set[str]) -> None:
    """docx / pptx 共通の段落置換。runを跨いで分断されたプレースホルダーは先頭runへ集約する。"""
    if "{{" not in p.text:
        return
    for run in p.runs:
        if "{{" in run.text:
            run.text = _apply_text(run.text, values, counter, missing)
    full = "".join(run.text for run in p.runs)
    replaced = _apply_text(full, values, counter, missing)
    if replaced != full and p.runs:
        p.runs[0].text = replaced
        for run in p.runs[1:]:
            run.text = ""


def _iter_docx_containers(doc):
    """本文・表(入れ子含む)・ヘッダー/フッターの段落コンテナを全て辿る。"""

    def walk(container):
        yield container
        for t in container.tables:
            for row in t.rows:
                for cell in row.cells:
                    yield from walk(cell)

    yield from walk(doc)
    for section in doc.sections:
        for part in (section.header, section.footer):
            yield from walk(part)


def _merge_docx(path: str, out_path, values, counter, missing) -> None:
    doc = Document(path)
    for container in _iter_docx_containers(doc):
        for p in container.paragraphs:
            _replace_in_paragraph(p, values, counter, missing)
    if counter:
        atomic_save(doc.save, out_path)


def _merge_xlsx(path: str, out_path, values, counter, missing) -> None:
    wb = load_workbook(path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                if not (isinstance(v, str) and "{{" in v):
                    continue
                m = PLACEHOLDER_RE.fullmatch(v.strip())
                if m and m.group(1) in values:
                    # セル全体が1つのプレースホルダーなら型を保って書き込む(数値は数値のまま)
                    counter[m.group(1)] = counter.get(m.group(1), 0) + 1
                    cell.value = values[m.group(1)]
                else:
                    cell.value = _apply_text(v, values, counter, missing)
    if counter:
        atomic_save(wb.save, out_path)


def _merge_pptx(path: str, out_path, values, counter, missing) -> None:
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    _replace_in_paragraph(p, values, counter, missing)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            _replace_in_paragraph(p, values, counter, missing)
    if counter:
        atomic_save(prs.save, out_path)


@tool
def merge_template(template_filename: str, output_filename: str, values: dict[str, Any]) -> str:
    """テンプレート文書内の {{キー}} プレースホルダーを値で置き換え、output_filenameとして保存する。
    Word/Excel/PowerPointすべてに使える。テンプレートは変更されずに残るため、
    請求書・報告書・案内文など同じ体裁の文書を宛先だけ変えて量産する用途に最適。
    テンプレートと同じファイル名を指定するとそのファイル自体を置換する。
    例: values={"宛先": "株式会社サンプル", "合計": 5200} → 文書内の {{宛先}} {{合計}} が置き換わる。"""
    src = resolve_workspace_path(template_filename, must_exist=True)
    out = resolve_workspace_path(output_filename)
    if src.suffix != out.suffix:
        return f"エラー: 出力ファイルはテンプレートと同じ拡張子({src.suffix})にしてください"
    if not values:
        return "エラー: valuesが空です。置き換えるキーと値を指定してください"

    counter: dict[str, int] = {}
    missing: set[str] = set()
    merger = {".docx": _merge_docx, ".xlsx": _merge_xlsx, ".pptx": _merge_pptx}.get(src.suffix)
    if merger is None:
        return "エラー: .docx / .xlsx / .pptx のいずれかを指定してください"
    merger(str(src), out, values, counter, missing)

    if not counter:
        found = "、".join(f"{{{{{k}}}}}" for k in sorted(missing)) if missing else "なし"
        return (
            f"エラー: 指定されたキーのプレースホルダーが1つも見つかりませんでした。"
            f"文書内にあるプレースホルダー: {found}"
        )
    total = sum(counter.values())
    detail = "、".join(f"{k}: {n}箇所" for k, n in counter.items())
    parts = [f"{output_filename} に差し込みました (計{total}箇所: {detail})"]
    unused = [k for k in values if k not in counter]
    if unused:
        parts.append(f"文書内に見つからなかったキー: {', '.join(unused)}")
    if missing:
        parts.append("値が未指定のまま残っている箇所: " + "、".join(f"{{{{{k}}}}}" for k in sorted(missing)))
    return "\n".join(parts)


TEMPLATE_TOOLS = [merge_template]
