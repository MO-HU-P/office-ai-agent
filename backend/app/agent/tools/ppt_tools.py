"""PowerPoint (.pptx) 操作ツール群"""
from pathlib import Path
from typing import Any, Optional

from langchain_core.tools import tool
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.text.text import _Run
from pptx.util import Cm, Pt

from ...atomic import atomic_save
from ...config import resolve_workspace_path
from .inline_format import describe_format, find_keyword_spans, split_runs_at_spans, validate_format_args

LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6

# 同梱のデザインテンプレート(16:9・アプリUIと同じGoogle風配色。assets/build_theme.py で生成)
_TEMPLATE_PATH = Path(__file__).parent.parent / "assets" / "default_theme.pptx"


def _open(filename: str):
    path = resolve_workspace_path(filename, must_exist=True)
    if path.suffix != ".pptx":
        raise ValueError("PowerPointは .pptx ファイルを指定してください")
    return Presentation(str(path)), str(path)


def _fit_width(prs, left: float, width: float) -> float:
    """既定の幅がスライド(4:3の既存ファイル等)からはみ出す場合、右余白0.5cmを残して縮める。"""
    slide_w = prs.slide_width / 360000  # cm
    if left + width > slide_w - 0.5:
        width = max(slide_w - left - 0.5, 2.0)
    return width


def _fill_bullets(body_shape, bullets: list[str]):
    tf = body_shape.text_frame
    tf.clear()
    first = True
    for item in bullets:
        stripped = item.lstrip(" ")
        level = min((len(item) - len(stripped)) // 2, 4)  # 先頭スペース2つでインデント1段
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = stripped
        p.level = level
        for run in p.runs:
            run.font.size = Pt(20 if level == 0 else 18)
        first = False


_TITLE_PH_TYPES = (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
_BODY_PH_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)


def _pick_layout(prs, preferred: int, need_body: bool):
    """スライドレイアウトを選ぶ。同梱テンプレートや標準的なファイルでは preferred のインデックスを
    そのまま使い、レイアウト構成が特殊なファイル(枚数が少ない・本文プレースホルダーが無い等)では、
    必要なプレースホルダーを実際に持つレイアウトを探して代わりに使う。"""

    def ph_types(layout):
        return [ph.placeholder_format.type for ph in layout.placeholders]

    def ok(layout):
        types = ph_types(layout)
        if not any(t in _TITLE_PH_TYPES for t in types):
            return False
        return (not need_body) or any(t in _BODY_PH_TYPES for t in types)

    layouts = list(prs.slide_layouts)
    if preferred < len(layouts) and ok(layouts[preferred]):
        return layouts[preferred]
    candidates = [l for l in layouts if ok(l)]
    if candidates:
        if not need_body:  # タイトルのみが目的なら、本文プレースホルダーの無いレイアウトを優先する
            no_body = [l for l in candidates if not any(t in _BODY_PH_TYPES for t in ph_types(l))]
            if no_body:
                return no_body[0]
        return candidates[0]
    # 条件を満たすレイアウトが無い→一番シンプルなものを使い、不足分はテキストボックスで補う
    return min(layouts, key=lambda l: len(l.placeholders)) if layouts else None


def _find_body_placeholder(slide):
    """箇条書きを入れられる本文プレースホルダーを探す(無ければNone)。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.type in _BODY_PH_TYPES and ph.has_text_frame:
            return ph
    return None


def _add_title_textbox(prs, slide, title: str):
    """タイトルプレースホルダーの無いスライドに、タイトル代わりのテキストボックスを置く。"""
    w = prs.slide_width / 360000
    h = prs.slide_height / 360000
    box = slide.shapes.add_textbox(Cm(w * 0.05), Cm(h * 0.05), Cm(w * 0.9), Cm(h * 0.14))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = title
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
    return box


def _title_bottom_cm(slide) -> Optional[float]:
    """タイトルに文字が入っている場合、その下端の位置(cm)を返す。タイトルが無い/空なら None。"""
    try:
        title = slide.shapes.title
        if title is None or not title.has_text_frame or not title.text_frame.text.strip():
            return None
        if title.top is None or title.height is None:
            return None
        return (title.top + title.height) / 360000
    except (KeyError, AttributeError, ValueError):
        return None


def _avoid_title_overlap(prs, slide, top: float) -> tuple[float, str]:
    """図形・表・グラフ・画像がタイトルに重なる位置に置かれそうなら、タイトルの下へ押し下げる。
    (調整後のtop, 報告用の注記) を返す。下げるとスライドに収まらない場合は動かさず警告だけ返す。"""
    bottom = _title_bottom_cm(slide)
    if bottom is None or top >= bottom + 0.2:
        return top, ""
    new_top = round(bottom + 0.3, 1)
    slide_h = prs.slide_height / 360000
    if new_top > slide_h - 3.0:
        return top, f"\n⚠️ タイトルと重なっている可能性があります(上端{top}cm)。render_pageやppt_read(mode=\"shapes\")で確認してください"
    return new_top, f"\n(タイトルとの重なりを避けるため、上端を{top}cm→{new_top}cmに自動調整しました)"


@tool
def ppt_create(filename: str, title: str = "", subtitle: str = "") -> str:
    """新しいPowerPoint(.pptx)を作成する。filenameは必ず.pptxで終わること。
    titleを指定するとタイトルスライドが追加される。
    16:9のデザイン済みテンプレート(配色・フォント・装飾入り)で作られる。"""
    path = resolve_workspace_path(filename)
    if path.suffix != ".pptx":
        return "エラー: filenameは .pptx で終わる必要があります"
    prs = Presentation(str(_TEMPLATE_PATH)) if _TEMPLATE_PATH.exists() else Presentation()
    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
    atomic_save(prs.save, path)
    return f"{filename} を作成しました"


@tool
def ppt_add_slide(filename: str, title: str, bullets: Optional[list[str]] = None, section_header: bool = False) -> str:
    """PowerPointにスライドを1枚追加する。bulletsは箇条書きのリスト。
    項目の先頭にスペース2つを付けると1段インデントされる(例: ["親項目", "  子項目"])。
    section_header=Trueにすると章の区切り用スライドになる。"""
    prs, path = _open(filename)
    if section_header:
        layout = _pick_layout(prs, LAYOUT_SECTION, need_body=False)
    elif bullets:
        layout = _pick_layout(prs, LAYOUT_TITLE_CONTENT, need_body=True)
    else:
        layout = _pick_layout(prs, LAYOUT_TITLE_ONLY, need_body=False)
    if layout is None:
        return "エラー: このファイルにはスライドレイアウトが1つもありません"
    slide = prs.slides.add_slide(layout)
    # レイアウトにタイトル・本文のプレースホルダーが無いファイルでは、テキストボックスで代用する
    notes = []
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        _add_title_textbox(prs, slide, title)
        notes.append("タイトルはテキストボックスで追加")
    if bullets:
        body = _find_body_placeholder(slide)
        if body is not None:
            _fill_bullets(body, bullets)
        else:
            w = prs.slide_width / 360000
            h = prs.slide_height / 360000
            box = slide.shapes.add_textbox(Cm(w * 0.06), Cm(h * 0.24), Cm(w * 0.88), Cm(h * 0.66))
            box.text_frame.word_wrap = True
            _fill_bullets(box, bullets)
            notes.append("本文はテキストボックスで追加")
    atomic_save(prs.save, path)
    result = f"{filename} にスライド{len(prs.slides)}「{title}」を追加しました"
    if notes:
        result += f"\n(このファイルには標準のプレースホルダーが無いため: {'、'.join(notes)})"
    return result


_SHAPE_TYPE_NAMES = {
    MSO_SHAPE_TYPE.PICTURE: "画像",
    MSO_SHAPE_TYPE.TEXT_BOX: "テキストボックス",
    MSO_SHAPE_TYPE.AUTO_SHAPE: "図形",
    MSO_SHAPE_TYPE.TABLE: "表",
    MSO_SHAPE_TYPE.CHART: "グラフ",
    MSO_SHAPE_TYPE.GROUP: "グループ",
    MSO_SHAPE_TYPE.LINE: "線",
}


def _shape_label(shape) -> str:
    if shape.is_placeholder:
        return "プレースホルダー(タイトル)" if shape.placeholder_format.idx == 0 else "プレースホルダー(本文)"
    try:
        return _SHAPE_TYPE_NAMES.get(shape.shape_type, str(shape.shape_type))
    except NotImplementedError:
        return "その他"


def _shape_line(num: int, shape) -> str:
    def cm(v):
        return f"{v.cm:.1f}" if v is not None else "?"

    line = (f"[{num}] {_shape_label(shape)}: 左{cm(shape.left)} 上{cm(shape.top)} "
            f"幅{cm(shape.width)} 高さ{cm(shape.height)}")
    if shape.has_text_frame:
        text = shape.text_frame.text.strip().replace("\n", " ")
        if text:
            line += f" 「{text[:25]}{'…' if len(text) > 25 else ''}」"
    return line


@tool
def ppt_read(filename: str, mode: str = "full") -> str:
    """PowerPointの内容を読む。スライド番号ごとにタイトルと本文テキストを返す。編集前に必ず呼ぶこと。
    mode="outline"にすると各スライドのタイトルだけを返す。枚数が多いときはまずoutlineで全体を把握すること。
    mode="shapes"にすると各スライドの図形一覧(図形番号・種類・位置と大きさcm)を返す。
    図形の削除(ppt_delete_shape)や移動(ppt_move_shape)の前に、この番号を確認すること。"""
    prs, _ = _open(filename)
    lines = [f"全{len(prs.slides)}枚"]
    if mode == "outline":
        for i, slide in enumerate(prs.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title is not None else ""
            lines.append(f"スライド{i}: {title or '(タイトルなし)'}")
        return "\n".join(lines)
    if mode == "shapes":
        lines[0] += f" (スライド寸法: 幅{prs.slide_width.cm:.1f}cm x 高さ{prs.slide_height.cm:.1f}cm)"
        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- スライド{i} ---")
            for num, shape in enumerate(slide.shapes, start=1):
                lines.append(_shape_line(num, shape))
        return "\n".join(lines)
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- スライド{i} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                text = "".join(run.text for run in p.runs)
                if text.strip():
                    prefix = "  " * p.level
                    lines.append(f"{prefix}{text}")
    return "\n".join(lines)


def _apply_slide_edit(prs, slide_number: int, title: str, bullets: Optional[list[str]]) -> Optional[str]:
    """1枚分の編集を適用する。失敗時はエラーメッセージを返し、成功時はNoneを返す。"""
    if not isinstance(slide_number, int) or slide_number < 1 or slide_number > len(prs.slides):
        return f"スライド番号 {slide_number} は範囲外です (1〜{len(prs.slides)})"
    slide = prs.slides[slide_number - 1]
    if title and slide.shapes.title is not None:
        slide.shapes.title.text = title
    if bullets is not None:
        body = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx != 0 and shape.has_text_frame:
                body = shape
                break
        if body is None:
            return f"スライド{slide_number}には本文プレースホルダーがありません"
        _fill_bullets(body, bullets)
    return None


@tool
def ppt_edit_slide(filename: str, slide_number: int, title: str = "", bullets: Optional[list[str]] = None) -> str:
    """既存スライドのタイトル・箇条書きを書き換える。slide_numberは1始まり(ppt_readの番号)。
    titleまたはbulletsのうち指定したものだけが置き換えられる。"""
    prs, path = _open(filename)
    error = _apply_slide_edit(prs, slide_number, title, bullets)
    if error:
        return f"エラー: {error}"
    atomic_save(prs.save, path)
    return f"スライド{slide_number}を更新しました"


@tool
def ppt_batch_edit(filename: str, edits: list[dict]) -> str:
    """複数のスライドをまとめて書き換える。2枚以上を直すときはppt_edit_slideを繰り返さず必ずこちらを使う。
    editsの各要素は {"slide_number": スライド番号(1始まり), "title": "新タイトル", "bullets": ["箇条書き", ...]}。
    titleとbulletsは指定したものだけが置き換えられる。一部が失敗しても残りは適用される。"""
    prs, path = _open(filename)
    ok_count = 0
    failures: list[str] = []
    for e in edits:
        error = _apply_slide_edit(prs, e.get("slide_number"), e.get("title") or "", e.get("bullets"))
        if error:
            failures.append(error)
        else:
            ok_count += 1
    if ok_count:
        atomic_save(prs.save, path)
    if not ok_count:
        return "エラー: 1枚も更新できませんでした\n" + "\n".join(failures)
    result = f"{filename} の{ok_count}枚のスライドを更新しました"
    if failures:
        result += "\n更新できなかったもの:\n" + "\n".join(failures)
    return result


_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_down": MSO_SHAPE.DOWN_ARROW,
    "star": MSO_SHAPE.STAR_5_POINT,
}


def _get_slide(prs, slide_number: int):
    """1始まりのスライド番号からスライドを取得する。範囲外ならNone。"""
    if not isinstance(slide_number, int) or slide_number < 1 or slide_number > len(prs.slides):
        return None
    return prs.slides[slide_number - 1]


def _slide_range_error(prs, slide_number) -> str:
    return f"エラー: スライド番号 {slide_number} は範囲外です (1〜{len(prs.slides)})"


@tool
def ppt_add_shape(
    filename: str,
    slide_number: int,
    shape: str = "rounded_rectangle",
    text: str = "",
    left: float = 2.0,
    top: float = 5.0,
    width: float = 8.0,
    height: float = 3.0,
    fill_color: str = "",
    font_size: float = 18,
    font_color: str = "",
) -> str:
    """既存スライドに図形またはテキストボックスを追加する。位置・大きさの単位はcm
    (既定の16:9スライドは幅33.87cm x 高さ19.05cm。実寸はppt_readで確認できる)。slide_numberは1始まり。
    shape: rectangle / rounded_rectangle / ellipse / arrow_right / arrow_down / star / text(枠なしテキストボックス)。
    textに改行(\\n)を入れると複数行になる。fill_color/font_colorは "#RRGGBB" 形式。"""
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    if shape != "text" and shape not in _SHAPE_MAP:
        return f"エラー: shapeは text / {' / '.join(_SHAPE_MAP)} のいずれかを指定してください"
    top, overlap_note = _avoid_title_overlap(prs, slide, top)
    if shape == "text":
        sp = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    else:
        sp = slide.shapes.add_shape(_SHAPE_MAP[shape], Cm(left), Cm(top), Cm(width), Cm(height))
        if fill_color:
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#").upper())
    tf = sp.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if shape != "text":
            p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(font_size)
            if font_color:
                run.font.color.rgb = RGBColor.from_string(font_color.lstrip("#").upper())
    atomic_save(prs.save, path)
    label = "テキストボックス" if shape == "text" else f"図形({shape})"
    return f"スライド{slide_number}に{label}を追加しました{overlap_note}"


@tool
def ppt_add_image(filename: str, slide_number: int, image_file: str, left: float = 2.0, top: float = 4.0, width: float = 0) -> str:
    """既存スライドにワークスペース内の画像ファイル(PNG/JPEGなど)を挿入する。
    位置・大きさの単位はcm(既定の16:9スライドは幅33.87cm x 高さ19.05cm)。slide_numberは1始まり。
    widthを指定すると縦横比を保ったままその幅に拡大縮小され、省略すると元の大きさで挿入される。"""
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    image_path = resolve_workspace_path(image_file, must_exist=True)
    if image_path.suffix.lower() not in (".png", ".jpg", ".jpeg", ".gif", ".bmp"):
        return "エラー: 画像は .png / .jpg / .jpeg / .gif / .bmp のファイルを指定してください"
    top, overlap_note = _avoid_title_overlap(prs, slide, top)
    kwargs = {"width": Cm(width)} if width and width > 0 else {}
    slide.shapes.add_picture(str(image_path), Cm(left), Cm(top), **kwargs)
    atomic_save(prs.save, path)
    return f"スライド{slide_number}に画像 {image_file} を挿入しました{overlap_note}"


@tool
def ppt_add_table(
    filename: str,
    slide_number: int,
    rows: list[list[Any]],
    left: float = 1.7,
    top: float = 4.5,
    width: float = 30.5,
    header: bool = True,
    font_size: float = 14,
) -> str:
    """既存スライドに表を追加する。rowsは2次元配列(行のリスト)で、セルは文字列・数値どちらでもよい。
    header=Trueなら1行目が見出し行になる。
    位置・大きさの単位はcm(既定の16:9スライドは幅33.87cm x 高さ19.05cm)。slide_numberは1始まり。"""
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    if not rows or not rows[0]:
        return "エラー: rowsが空です"
    top, overlap_note = _avoid_title_overlap(prs, slide, top)
    width = _fit_width(prs, left, width)
    n_cols = max(len(r) for r in rows)
    height = Cm(min(1.2 * len(rows), 14.0))  # 行数に応じた高さ(あふれた分は自動で伸びる)
    table = slide.shapes.add_table(len(rows), n_cols, Cm(left), Cm(top), Cm(width), height).table
    table.first_row = header
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.text = str(row[ci]) if ci < len(row) and row[ci] is not None else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
    atomic_save(prs.save, path)
    return f"スライド{slide_number}に {len(rows)}x{n_cols} の表を追加しました{overlap_note}"


@tool
def ppt_add_chart(
    filename: str,
    slide_number: int,
    chart_type: str,
    categories: list[str],
    series: dict[str, list[float]],
    title: str = "",
    left: float = 1.7,
    top: float = 3.9,
    width: float = 30.5,
    height: float = 13.5,
) -> str:
    """既存スライドにPowerPoint上で編集できるグラフを追加する。画像ではなくネイティブのグラフなので、
    あとからPowerPointで数値や色を変更できる。位置・大きさの単位はcm。slide_numberは1始まり。
    chart_type: bar(縦棒) / bar_horizontal(横棒) / line(折れ線) / pie(円)。
    categoriesは項目名のリスト、seriesは系列名→数値リストの辞書。
    例: categories=["4月","5月"], series={"売上": [100, 120], "利益": [20, 30]}。
    pieは最初の1系列だけが使われる。"""
    type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    if chart_type not in type_map:
        return f"エラー: chart_typeは {' / '.join(type_map)} のいずれかを指定してください"
    if not categories or not series:
        return "エラー: categoriesとseriesの両方を指定してください"
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    top, overlap_note = _avoid_title_overlap(prs, slide, top)
    width = _fit_width(prs, left, width)
    data = CategoryChartData()
    data.categories = categories
    items = list(series.items())[:1] if chart_type == "pie" else series.items()
    for name, values in items:
        if len(values) != len(categories):
            return f"エラー: 系列「{name}」の数値の個数({len(values)})がcategoriesの個数({len(categories)})と一致しません"
        data.add_series(str(name), [float(v) for v in values])
    frame = slide.shapes.add_chart(type_map[chart_type], Cm(left), Cm(top), Cm(width), Cm(height), data)
    chart = frame.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    if chart_type == "pie" or len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    atomic_save(prs.save, path)
    return f"スライド{slide_number}に{chart_type}グラフを追加しました{overlap_note}"


def _iter_text_frames(shapes):
    """図形の集まりから、文字の入る全テキストフレームを返す(グループの中・表のセルも含む)。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_frames(shape.shapes)
            continue
        if shape.has_text_frame:
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


@tool
def ppt_format_text(
    filename: str,
    keywords: list[str],
    slide_number: int = 0,
    color: str = "",
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
) -> str:
    """スライド内の特定の語句(キーワード)だけに文字書式を付ける。「重要な用語を赤字にして」
    「◯◯を太字にして」のような、文中の一部の文字だけの色付け・強調に使う。
    タイトル・本文・図形・表の中からkeywordsの全出現箇所を探して同じ書式を適用する。
    slide_numberは1始まりで、0(既定)なら全スライドが対象。colorは"#RRGGBB"形式(赤字なら"#FF0000")。
    bold/italic/underlineはtrueで付け、falseで外す。指定したものだけが変更される。"""
    keywords, error = validate_format_args(keywords, color, bold, italic, underline)
    if error:
        return error
    prs, path = _open(filename)
    if slide_number:
        slide = _get_slide(prs, slide_number)
        if slide is None:
            return _slide_range_error(prs, slide_number)
        slides = [slide]
    else:
        slides = list(prs.slides)
    rgb = RGBColor.from_string(color.lstrip("#").upper()) if color else None
    hit_count = 0
    for slide in slides:
        for tf in _iter_text_frames(slide.shapes):
            for p in tf.paragraphs:
                runs = list(p.runs)
                spans = find_keyword_spans("".join(r.text for r in runs), keywords)
                if not spans:
                    continue
                for el in split_runs_at_spans([r._r for r in runs], spans):
                    run = _Run(el, p)  # キーワード部分のrunにだけ書式を適用する
                    if rgb is not None:
                        run.font.color.rgb = rgb
                    if bold is not None:
                        run.font.bold = bold
                    if italic is not None:
                        run.font.italic = italic
                    if underline is not None:
                        run.font.underline = underline
                hit_count += len(spans)
    if not hit_count:
        where = f"スライド{slide_number}" if slide_number else filename
        return f"「{'」「'.join(keywords)}」は {where} に見つかりませんでした。ppt_readで実際の表記を確認してください"
    atomic_save(prs.save, path)
    where = f"スライド{slide_number}" if slide_number else f"{filename} 全体"
    return f"{where}の{hit_count}箇所に書式({describe_format(color, bold, italic, underline)})を適用しました"


@tool
def ppt_delete_shape(filename: str, slide_number: int, shape_numbers: list[int]) -> str:
    """スライドから図形・画像・テキストボックスなどを削除する。はみ出しや重なりを直すときに、
    上に別の図形をかぶせて隠すのではなく、不要な図形をこれで削除してから作り直すこと。
    shape_numbersは削除する図形番号のリスト(1始まり。必ず直前に ppt_read の mode="shapes" で確認する)。"""
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    shapes = list(slide.shapes)
    targets = []
    failures: list[str] = []
    for n in dict.fromkeys(shape_numbers):  # 重複指定は1回にまとめる
        if not isinstance(n, int) or n < 1 or n > len(shapes):
            failures.append(f"図形番号 {n} は範囲外です (1〜{len(shapes)})")
        else:
            targets.append((n, shapes[n - 1]))
    # 先に対象の要素を確定してから消すので、削除による番号ずれの影響を受けない
    for _, shape in targets:
        el = shape._element
        el.getparent().remove(el)
    if not targets:
        return "エラー: 1つも削除できませんでした\n" + "\n".join(failures)
    atomic_save(prs.save, path)
    nums = ", ".join(str(n) for n, _ in targets)
    result = f"スライド{slide_number}の図形 [{nums}] を削除しました(残り{len(shapes) - len(targets)}個)。番号が振り直されるため、続けて操作するときはppt_readのmode=\"shapes\"で確認し直すこと"
    if failures:
        result += "\n削除できなかったもの:\n" + "\n".join(failures)
    return result


@tool
def ppt_move_shape(
    filename: str,
    slide_number: int,
    shape_number: int,
    left: Optional[float] = None,
    top: Optional[float] = None,
    width: Optional[float] = None,
    height: Optional[float] = None,
) -> str:
    """スライド上の図形・画像・プレースホルダーの位置や大きさを変更する。はみ出し・重なりの修正に使う。
    shape_numberは図形番号(1始まり。ppt_read の mode="shapes" で確認する)。
    left/top/width/heightはcm単位で、指定したものだけが変更される(例: width=10 なら幅だけ変わる)。"""
    prs, path = _open(filename)
    slide = _get_slide(prs, slide_number)
    if slide is None:
        return _slide_range_error(prs, slide_number)
    shapes = list(slide.shapes)
    if not isinstance(shape_number, int) or shape_number < 1 or shape_number > len(shapes):
        return f"エラー: 図形番号 {shape_number} は範囲外です (1〜{len(shapes)})"
    if left is None and top is None and width is None and height is None:
        return "エラー: left / top / width / height のうち少なくとも1つを指定してください"
    shape = shapes[shape_number - 1]
    changed = []
    if left is not None:
        shape.left = Cm(left)
        changed.append(f"左{left}cm")
    if top is not None:
        shape.top = Cm(top)
        changed.append(f"上{top}cm")
    if width is not None:
        shape.width = Cm(width)
        changed.append(f"幅{width}cm")
    if height is not None:
        shape.height = Cm(height)
        changed.append(f"高さ{height}cm")
    atomic_save(prs.save, path)
    return f"スライド{slide_number}の図形 [{shape_number}] を変更しました ({', '.join(changed)})"


PPT_TOOLS = [
    ppt_create,
    ppt_add_slide,
    ppt_read,
    ppt_edit_slide,
    ppt_batch_edit,
    ppt_format_text,
    ppt_add_shape,
    ppt_add_image,
    ppt_add_table,
    ppt_add_chart,
    ppt_delete_shape,
    ppt_move_shape,
]
