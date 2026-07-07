"""個人情報(PII)のローカル・マスク処理。

LLMに生の個人情報を渡さず、Python側の正規表現だけで決定論的にマスクする。
対象は「形式が決まった構造化PII」= メールアドレス・電話番号・URL・郵便番号・
マイナンバー・クレジットカード番号。氏名・住所・会社名などパターンの無い情報は
正規表現では確実に検出できないため対象外(README.ja.md にその旨を明記)。
"""
import re
from pathlib import Path

from docx import Document
from langchain_core.tools import tool
from openpyxl import load_workbook
from pptx import Presentation

from ...atomic import atomic_save
from ...config import resolve_workspace_path

# URLの終端とみなす文字(空白・引用符・山括弧に加え、日本語の句読点・括弧・全角空白)。
# これらを含めないと「https://example.org、〒530-0001」のようにURLの後ろの
# 日本語(読点や郵便番号)まで巻き込んでしまう。
_URL_STOP = r'\s<>"　、。，．・…‥（）「」『』【】〈〉《》〔〕！？'

# 適用順が重要。特殊文字を含むURL・メールを先に消し、続いて桁数の多い順に数字列を消す。
# (先にマッチした部分は【…】に置換され、後続パターンの誤検出を防ぐ)
_PATTERNS: list[tuple[str, re.Pattern]] = [
    ("URL", re.compile(rf"(?:https?://|www\.)[^{_URL_STOP}]+")),
    ("メールアドレス", re.compile(r"[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}")),
    ("カード番号", re.compile(r"(?<!\d)\d{4}[-\s]\d{4}[-\s]\d{4}[-\s]\d{4}(?!\d)")),
    ("電話番号", re.compile(
        r"(?<!\d)(?:\+81[-\d]{9,13}"
        r"|0[5789]0-\d{4}-\d{4}"
        r"|0\d{1,4}-\d{1,4}-\d{4}"
        r"|\(0\d{1,4}\)\s?\d{1,4}-\d{4})(?!\d)")),
    ("郵便番号", re.compile(r"〒\s?\d{3}-\d{4}|(?<!\d)\d{3}-\d{4}(?!\d)")),
    ("マイナンバー", re.compile(r"(?<!\d)\d{12}(?!\d)")),
]


def mask_pii_text(text: str) -> tuple[str, dict[str, int]]:
    """文字列中の構造化PIIを【メールアドレス】等のプレースホルダに置換し、
    (置換後の文字列, 種類ごとの件数) を返す。"""
    counts: dict[str, int] = {}
    for label, pattern in _PATTERNS:

        def repl(_m, _label=label):
            counts[_label] = counts.get(_label, 0) + 1
            return f"【{_label}】"

        text = pattern.sub(repl, text)
    return text, counts


def _merge(dst: dict[str, int], src: dict[str, int]) -> None:
    for k, v in src.items():
        dst[k] = dst.get(k, 0) + v


def _mask_paragraph(p, counts: dict[str, int]) -> None:
    """段落(docx/pptx共通)のテキストをマスクする。ランをまたぐPIIも確実に消すため、
    段落テキスト全体でマッチさせ、変化があれば先頭ランに書式を残して1ランに束ねる。"""
    runs = list(p.runs)
    original = "".join(r.text for r in runs)
    if not original:
        return
    masked, c = mask_pii_text(original)
    if masked == original:
        return
    _merge(counts, c)
    if runs:
        runs[0].text = masked
        for r in runs[1:]:
            r._r.getparent().remove(r._r)


def _anonymize_docx(src: str, out: str) -> dict[str, int]:
    doc = Document(src)
    counts: dict[str, int] = {}
    for p in doc.paragraphs:
        _mask_paragraph(p, counts)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _mask_paragraph(p, counts)
    atomic_save(doc.save, Path(out))
    return counts


def _anonymize_xlsx(src: str, out: str) -> dict[str, int]:
    wb = load_workbook(src)
    counts: dict[str, int] = {}
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    masked, c = mask_pii_text(cell.value)
                    if masked != cell.value:
                        cell.value = masked
                        _merge(counts, c)
    atomic_save(wb.save, Path(out))
    return counts


def _anonymize_pptx(src: str, out: str) -> dict[str, int]:
    prs = Presentation(src)
    counts: dict[str, int] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    _mask_paragraph(p, counts)
    atomic_save(prs.save, Path(out))
    return counts


def _anonymize_text(src: str, out: str) -> dict[str, int]:
    original = Path(src).read_text(encoding="utf-8", errors="replace")
    masked, counts = mask_pii_text(original)
    atomic_save(lambda pth: Path(pth).write_text(masked, encoding="utf-8"), Path(out))
    return counts


_ANONYMIZERS = {
    ".docx": _anonymize_docx,
    ".xlsx": _anonymize_xlsx,
    ".pptx": _anonymize_pptx,
    ".csv": _anonymize_text,
    ".txt": _anonymize_text,
}


@tool
def anonymize_file(filename: str, dest: str = "") -> str:
    """ファイル内の個人情報を、元ファイルを残したまま自動でマスクしたコピーを作る。
    マスク対象はメールアドレス・電話番号・URL・郵便番号・マイナンバー・クレジットカード番号
    (形式が決まった情報)で、これらはAIに送らずローカルで確実に置き換える。
    ※氏名・住所・会社/学校名などパターンの無い情報は自動マスクできないので、その旨を必ず伝えること。
    対応形式は .docx / .xlsx / .pptx / .csv / .txt。destを省略すると「(匿名化)元の名前」で保存する。
    「匿名化して」「個人情報を伏せて」と言われたらこれを使う(自分で本文を書き換えない)。"""
    src = resolve_workspace_path(filename, must_exist=True)
    suffix = src.suffix.lower()
    anonymizer = _ANONYMIZERS.get(suffix)
    if anonymizer is None:
        return f"エラー: {suffix} は匿名化に対応していません(.docx / .xlsx / .pptx / .csv / .txt)"

    if dest:
        out = resolve_workspace_path(dest)
        if out.suffix.lower() != suffix:
            return f"エラー: destは元と同じ拡張子({suffix})にしてください"
    else:
        out = resolve_workspace_path(f"(匿名化){src.name}")
    if out.exists():
        return f"エラー: {out.name} は既に存在します。destで別の名前を指定してください"

    counts = anonymizer(str(src), str(out))
    note = "氏名・住所・会社/学校名などは自動マスクの対象外です。必要なら目視で確認してください。"
    total = sum(counts.values())
    if total == 0:
        # コピー自体は作らない(マスク対象が無いのにコピーだけ増えるのを避ける)
        if out.exists():
            out.unlink()
        return f"構造化された個人情報(メール・電話番号・URL・郵便番号・マイナンバー・カード番号)は見つかりませんでした。{note}"
    breakdown = "、".join(f"{k} {v}件" for k, v in counts.items())
    return f"{out.name} を作成し、{breakdown}をマスクしました(元の {filename} はそのまま残しています)。{note}"


PII_TOOLS = [anonymize_file]
