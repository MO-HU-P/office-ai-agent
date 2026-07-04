"""文書セルフチェックツール: 仕上がりの問題点をレンダリングなしで検出する"""
import re

from docx import Document
from langchain_core.tools import tool
from openpyxl import load_workbook
from pptx import Presentation

from ...config import resolve_workspace_path
from .template_tools import PLACEHOLDER_RE, _iter_docx_containers

# Excelの数式エラー値(キャッシュされた計算結果に現れる)
EXCEL_ERROR_VALUES = {"#REF!", "#DIV/0!", "#VALUE!", "#NAME?", "#N/A", "#NULL!", "#NUM!"}
# 数式中のシート参照: 'シート名'! または SheetName!
SHEET_REF_RE = re.compile(r"(?:'([^']+)'|([A-Za-z0-9_.]+))!")

MAX_ISSUES = 30


def _check_docx(path: str) -> list[str]:
    doc = Document(path)
    issues: list[str] = []
    empty_streak = 0
    for i, p in enumerate(doc.paragraphs):
        style = p.style.name if p.style else "Normal"
        text = p.text.strip()
        if (style == "Title" or style.startswith("Heading")) and not text:
            issues.append(f"段落[{i}]: 見出しスタイルなのに本文が空です")
        if not text:
            empty_streak += 1
            if empty_streak == 3:
                issues.append(f"段落[{i}]付近: 空行が3つ以上続いています(レイアウト崩れの可能性)")
        else:
            empty_streak = 0
    for container in _iter_docx_containers(doc):
        for p in container.paragraphs:
            for m in PLACEHOLDER_RE.finditer(p.text):
                issues.append(f"未置換のプレースホルダーが残っています: {{{{{m.group(1)}}}}}")
    if not any(p.text.strip() for p in doc.paragraphs) and not doc.tables:
        issues.append("文書に本文がありません")
    return issues


def _check_xlsx(path: str) -> list[str]:
    issues: list[str] = []
    wb = load_workbook(path)
    # data_only=True でキャッシュされた計算結果を読み、エラー値を検出する
    # (このアプリで書き込んだ直後の数式は未計算のため検出できない場合がある)
    wbv = load_workbook(path, data_only=True)
    sheet_names = set(wb.sheetnames)
    for ws, wsv in zip(wb.worksheets, wbv.worksheets):
        for row, row_v in zip(ws.iter_rows(), wsv.iter_rows()):
            for cell, cell_v in zip(row, row_v):
                v = cell.value
                if isinstance(cell_v.value, str) and cell_v.value in EXCEL_ERROR_VALUES:
                    issues.append(f"{ws.title}!{cell.coordinate}: 数式がエラー({cell_v.value})になっています")
                if isinstance(v, str):
                    if v.startswith("="):
                        for m in SHEET_REF_RE.finditer(v):
                            ref = m.group(1) or m.group(2)
                            if ref not in sheet_names:
                                issues.append(f"{ws.title}!{cell.coordinate}: 存在しないシート「{ref}」を参照しています")
                    else:
                        if len(v) > 80:
                            issues.append(f"{ws.title}!{cell.coordinate}: 文字列が長すぎます({len(v)}文字)。セル内改行や列分割を検討")
                        for m in PLACEHOLDER_RE.finditer(v):
                            issues.append(f"{ws.title}!{cell.coordinate}: 未置換のプレースホルダー {{{{{m.group(1)}}}}}")
    return issues


def _check_pptx(path: str) -> list[str]:
    issues: list[str] = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        if not title:
            issues.append(f"スライド{i}: タイトルがありません")
        elif len(title) > 40:
            issues.append(f"スライド{i}: タイトルが長すぎます({len(title)}文字)。はみ出す可能性があります")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if shape.is_placeholder and shape.placeholder_format.idx != 0 and not paragraphs:
                issues.append(f"スライド{i}: 未入力のプレースホルダーがあります")
            if len(paragraphs) > 8:
                issues.append(f"スライド{i}: 箇条書きが{len(paragraphs)}行あります。8行以下に分割を検討")
            for p in shape.text_frame.paragraphs:
                for m in PLACEHOLDER_RE.finditer(p.text):
                    issues.append(f"スライド{i}: 未置換のプレースホルダー {{{{{m.group(1)}}}}}")
    return issues


@tool
def check_document_issues(filename: str) -> str:
    """文書の仕上がりをチェックし、直したほうがよい点を列挙する。Word/Excel/PowerPointすべてに使える。
    検出内容: 空の見出し・連続空行(Word)、数式エラー・存在しないシート参照・長すぎる文字列(Excel)、
    タイトル欠落・箇条書き過多(PowerPoint)、未置換の {{プレースホルダー}}(共通)。
    文書を作り終えた後の最終確認や、ユーザーに「確認して」と言われたときに使う。"""
    path = resolve_workspace_path(filename, must_exist=True)
    checker = {".docx": _check_docx, ".xlsx": _check_xlsx, ".pptx": _check_pptx}.get(path.suffix)
    if checker is None:
        return "エラー: .docx / .xlsx / .pptx のいずれかを指定してください"
    issues = checker(str(path))
    if not issues:
        return f"{filename}: 問題は見つかりませんでした"
    shown = issues[:MAX_ISSUES]
    lines = [f"{filename}: {len(issues)}件の気になる点があります"] + [f"・{s}" for s in shown]
    if len(issues) > MAX_ISSUES:
        lines.append(f"…ほか{len(issues) - MAX_ISSUES}件")
    return "\n".join(lines)


CHECK_TOOLS = [check_document_issues]
