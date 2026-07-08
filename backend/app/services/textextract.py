"""Office/テキストファイルから「比較用の行」を取り出す共通処理。

doc_diff(2ファイル比較ツール)と変更差分ビュー(/api/files/{name}/changes)の両方で使う。
書式ではなく中身(文章・セルの値)を行のリストにする。
"""
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def docx_lines(path: str) -> list[str]:
    doc = Document(path)
    lines = [p.text for p in doc.paragraphs]
    for ti, table in enumerate(doc.tables):
        lines.append(f"[表{ti}]")
        for row in table.rows:
            lines.append(" | ".join(c.text for c in row.cells))
    return lines


def xlsx_lines(path: str) -> list[str]:
    """行単位の抽出(2ファイル比較用)。1行のセルを « | » で連結する。"""
    wb = load_workbook(path, read_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"=== シート: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            lines.append(" | ".join(cells).rstrip(" |"))
    wb.close()
    return lines


def xlsx_cell_lines(path: str) -> list[str]:
    """セル単位の抽出(変更差分ビュー用)。「シート名!A1: 値」の形式で、
    どのセルが変わったかをセル番地つきで示せる。空セルは出さない。"""
    wb = load_workbook(path, read_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                lines.append(f"{ws.title}!{cell.coordinate}: {cell.value}")
    wb.close()
    return lines


def pptx_lines(path: str) -> list[str]:
    prs = Presentation(path)
    lines: list[str] = []
    for si, slide in enumerate(prs.slides, 1):
        lines.append(f"=== スライド{si} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    lines.append("".join(run.text for run in para.runs))
    return lines


def text_lines(path: str) -> list[str]:
    return Path(path).read_text(encoding="utf-8", errors="replace").splitlines()


# 2ファイル比較(doc_diff)用: 行単位
EXTRACTORS = {
    ".docx": docx_lines,
    ".xlsx": xlsx_lines,
    ".pptx": pptx_lines,
    ".csv": text_lines,
    ".txt": text_lines,
}

# 変更差分ビュー用: Excelだけセル単位に置き換える(どのセルが変わったかを見せるため)
CHANGE_VIEW_EXTRACTORS = {**EXTRACTORS, ".xlsx": xlsx_cell_lines}
