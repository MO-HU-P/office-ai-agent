"""右ペイン用のプレビュー生成と、エージェント用のページ画像化。

- Excel: openpyxlでセル値+基本スタイルをJSON化(フロントでSheets風グリッド描画)。
  シートに貼られた画像(グラフのPNGなど)はキャッシュに書き出し、フロントがグリッドに重ねる
- PowerPoint: LibreOffice headlessでPDF化 → pdftoppmでスライドPNG化
- Word: フロント側で docx-preview がrawファイルを直接描画するため変換不要
  (エージェントのrender_pageツールはWordもLibreOfficeでPNG化する)
"""
import asyncio
import datetime as dt
import glob as glob_module
import shutil
import subprocess
import threading
from pathlib import Path
from typing import Any

from openpyxl import load_workbook
from openpyxl.utils import get_column_letter, range_to_tuple
from pptx import Presentation

from ..config import PREVIEW_CACHE_DIR

MAX_ROWS = 300
MAX_COLS = 60

# Office内部の長さの単位。96dpi換算でピクセルに直す
EMU_PER_PX = 9525

# LibreOfficeは同時起動に弱いため直列化。プレビューAPI(イベントループ)と
# エージェントツール(ワーカースレッド)の両方から呼ばれるのでthreading.Lockを使う
_soffice_lock = threading.Lock()

# Excelの画像をキャッシュへ書き出す処理の排他(同じファイルへの同時リクエスト対策)
_excel_img_lock = threading.Lock()


def _color_hex(color) -> str | None:
    try:
        rgb = color.rgb
    except AttributeError:
        return None
    if not isinstance(rgb, str) or len(rgb) != 8:
        return None
    if rgb == "00000000":
        return None
    return "#" + rgb[2:]


def _display_value(value, number_format: str) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, dt.datetime):
        if value.hour == 0 and value.minute == 0 and value.second == 0:
            return value.strftime("%Y/%m/%d")
        return value.strftime("%Y/%m/%d %H:%M")
    if isinstance(value, dt.date):
        return value.strftime("%Y/%m/%d")
    if isinstance(value, dt.time):
        return value.strftime("%H:%M")
    if isinstance(value, float):
        nf = number_format or ""
        if "%" in nf:
            digits = 1 if "0.0" in nf else 0
            return f"{value * 100:.{digits}f}%"
        if value == int(value) and abs(value) < 1e15:
            value = int(value)
        else:
            value = round(value, 6)
    if isinstance(value, (int, float)) and "#,##" in (number_format or ""):
        return f"{value:,}"
    return str(value)


def _image_size(img) -> tuple[int, int]:
    """画像の表示サイズ(px)。アンカーにextがあればそれ(=Excel上の実サイズ)、無ければ元画像のサイズ。"""
    ext = getattr(img.anchor, "ext", None)
    if ext is not None:
        return max(round(ext.cx / EMU_PER_PX), 1), max(round(ext.cy / EMU_PER_PX), 1)
    return max(int(img.width or 1), 1), max(int(img.height or 1), 1)


def _excel_images(wb, path: Path) -> dict[str, list[dict[str, Any]]]:
    """各シートに貼られた画像を取り出し、シート名 → 画像情報(位置・サイズ・URL)を返す。

    Excelプレビューはセルを自前描画していて画像を持たないため、フロントでグリッドに重ねる。
    画像の実体はPPTXのスライドPNGと同じmtimeキーのキャッシュに書き出し、URLで渡す。
    """
    by_sheet: dict[str, list[dict[str, Any]]] = {}
    pending: list[tuple[str, bytes]] = []
    cache_name = _cache_dir_for(path).name
    for si, name in enumerate(wb.sheetnames):
        items: list[dict[str, Any]] = []
        for i, img in enumerate(getattr(wb[name], "_images", [])):
            # セル基準でないアンカー(ページ絶対座標)はグリッド上に置けないので出さない
            frm = getattr(img.anchor, "_from", None)
            if frm is None:
                continue
            try:
                data = img._data()
            except Exception:
                continue  # 壊れた画像でプレビュー全体を落とさない
            # 拡張子はファイル名とURLに入るので、想定した形式以外はpng扱いにする
            fmt = (getattr(img, "format", None) or "png").lower()
            if fmt not in ("png", "jpg", "jpeg", "gif", "bmp"):
                fmt = "png"
            fname = f"xlimg-{si}-{i}.{fmt}"
            w, h = _image_size(img)
            items.append({
                "url": f"/api/preview_cache/{cache_name}/{fname}",
                "r": frm.row,
                "c": frm.col,
                "dx": round((frm.colOff or 0) / EMU_PER_PX),
                "dy": round((frm.rowOff or 0) / EMU_PER_PX),
                "w": w,
                "h": h,
            })
            pending.append((fname, data))
        if items:
            by_sheet[name] = items
    if not pending:
        return {}
    cache_dir = _cache_dir_for(path)
    with _excel_img_lock:
        if not cache_dir.exists():
            tmp_dir = cache_dir.with_suffix(".tmp")
            shutil.rmtree(tmp_dir, ignore_errors=True)
            tmp_dir.mkdir(parents=True)
            try:
                for fname, data in pending:
                    (tmp_dir / fname).write_bytes(data)
                # 同じファイルの古いキャッシュを掃除してから確定(_convert_to_pngsと同じ手順)
                for old in PREVIEW_CACHE_DIR.glob(f"{glob_module.escape(path.stem)}__*"):
                    if old != tmp_dir:
                        shutil.rmtree(old, ignore_errors=True)
                tmp_dir.rename(cache_dir)
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
    return by_sheet


def _chart_size(chart) -> tuple[int, int]:
    """グラフの表示サイズ(px)。アンカーのextが基本。無ければopenpyxlの既定であるcm指定から換算する。"""
    ext = getattr(chart.anchor, "ext", None)
    if ext is not None:
        return max(round(ext.cx / EMU_PER_PX), 1), max(round(ext.cy / EMU_PER_PX), 1)
    return round((chart.width or 15) / 2.54 * 96), round((chart.height or 7.5) / 2.54 * 96)


def _rich_text(title) -> str:
    """グラフ・軸のタイトルから文字列を取り出す(取れなければ空文字)。"""
    if title is None:
        return ""
    try:
        rich = title.tx.rich
        if rich is not None:
            return "".join(r.t or "" for p in rich.p for r in (p.r or [])).strip()
    except AttributeError:
        pass
    return ""


def _ref_values(wb, ref: str | None) -> list[Any]:
    """"'シート'!$B$2:$B$4" のような参照を、その範囲のセルの値のリストに解決する。"""
    if not ref:
        return []
    try:
        sheet_name, (min_col, min_row, max_col, max_row) = range_to_tuple(ref)
    except Exception:
        return []  # 別ファイル参照や名前付き範囲などは解決しない
    if sheet_name not in wb.sheetnames:
        return []
    ws = wb[sheet_name]
    out = []
    for row in range(min_row, max_row + 1):
        for col in range(min_col, max_col + 1):
            out.append(ws.cell(row=row, column=col).value)
    return out


def _series_color(s) -> str | None:
    """系列に色が明示されていればその色(#rrggbb)。無ければNone(フロント側の既定色を使う)。"""
    fill = getattr(getattr(s, "graphicalProperties", None), "solidFill", None)
    if fill is None:
        return None
    rgb = fill if isinstance(fill, str) else getattr(fill, "srgbClr", None)
    # srgbClrはColorChoice内でRGB型に包まれていることがある
    rgb = getattr(rgb, "value", rgb)
    if isinstance(rgb, str) and len(rgb) == 6:
        return f"#{rgb.lower()}"
    return None


def _error_bars(wb, s, n: int) -> list[float | None] | None:
    """系列のエラーバー(誤差範囲)の値。セル参照で指定されたもの(errValType='cust')だけ扱う。"""
    eb = getattr(s, "errBars", None)
    if eb is None or getattr(eb, "errValType", None) != "cust":
        return None
    ref = getattr(getattr(eb, "plus", None), "numRef", None)
    if ref is None or not getattr(ref, "f", None):
        return None
    values = _ref_values(wb, ref.f)
    if not any(isinstance(v, (int, float)) and not isinstance(v, bool) for v in values):
        return None
    out: list[float | None] = [None] * n
    for i, v in enumerate(values[:n]):
        out[i] = float(v) if isinstance(v, (int, float)) and not isinstance(v, bool) else None
    return out


def _point_colors(s, n: int) -> list[str | None]:
    """円グラフ用。データ点ごとに色が指定されていればその色を、無ければNoneを並べて返す。"""
    out: list[str | None] = [None] * n
    for dp in getattr(s, "data_points", None) or []:
        idx = getattr(dp, "idx", None)
        if idx is None or not (0 <= idx < n):
            continue
        out[idx] = _series_color(dp)
    return out


def _series_ref(part) -> str | None:
    """系列の値/項目名/系列名が指すセル参照の文字列を返す。"""
    for kind in ("numRef", "strRef"):
        ref = getattr(part, kind, None)
        if ref is not None and getattr(ref, "f", None):
            return ref.f
    return None


# openpyxlのグラフ種別 → フロントが描けるグラフの種類
def _chart_kind(chart) -> str | None:
    tag = getattr(chart, "tagname", "")
    if tag == "barChart":
        return "bar_horizontal" if getattr(chart, "type", "col") == "bar" else "bar"
    if tag == "lineChart":
        return "line"
    if tag in ("pieChart", "doughnutChart"):
        return "pie"
    return None  # 散布図・面グラフなどは未対応(注記だけ出す)


def _excel_charts(wb, wb_f) -> dict[str, list[dict[str, Any]]]:
    """各シートのネイティブグラフを、フロントがSVGで描ける形の定義にして返す。

    値は参照先のセルから読むので、セルを直せばプレビューのグラフも追従する
    (Excelでの見え方と同じ)。wb_fはグラフ定義の取得用(data_only=Trueだと欠けるため)。
    """
    by_sheet: dict[str, list[dict[str, Any]]] = {}
    for name in wb_f.sheetnames:
        items: list[dict[str, Any]] = []
        for chart in getattr(wb_f[name], "_charts", []):
            frm = getattr(chart.anchor, "_from", None)
            if frm is None:
                continue
            w, h = _chart_size(chart)
            spec: dict[str, Any] = {
                "r": frm.row,
                "c": frm.col,
                "dx": round((frm.colOff or 0) / EMU_PER_PX),
                "dy": round((frm.rowOff or 0) / EMU_PER_PX),
                "w": w,
                "h": h,
                "title": _rich_text(chart.title),
                "kind": _chart_kind(chart),
            }
            if spec["kind"] is None:
                items.append(spec)  # 未対応の種類は「グラフがあります」の注記だけ出す
                continue
            try:
                spec["xTitle"] = _rich_text(chart.x_axis.title)
                spec["yTitle"] = _rich_text(chart.y_axis.title)
            except AttributeError:
                spec["xTitle"] = spec["yTitle"] = ""
            spec["stacked"] = getattr(chart, "grouping", "") in ("stacked", "percentStacked")
            categories: list[str] = []
            series: list[dict[str, Any]] = []
            for s in chart.series:
                values = [v if isinstance(v, (int, float)) and not isinstance(v, bool) else None
                          for v in _ref_values(wb, _series_ref(s.val))]
                if not any(v is not None for v in values):
                    continue
                label = _ref_values(wb, _series_ref(s.tx)) if s.tx is not None else []
                series.append({
                    "name": str(label[0]) if label and label[0] is not None else "",
                    "values": values,
                    "color": _series_color(s),
                    "errors": _error_bars(wb, s, len(values)),
                })
                if not categories and s.cat is not None:
                    categories = [("" if v is None else str(v)) for v in _ref_values(wb, _series_ref(s.cat))]
            if not series:
                spec["kind"] = None  # 値が読めないグラフは注記に落とす
                items.append(spec)
                continue
            # 項目名が取れないときは1,2,3…で補う(Excelも同じ扱い)
            n = max(len(s["values"]) for s in series)
            if len(categories) < n:
                categories += [str(i + 1) for i in range(len(categories), n)]
            spec["categories"] = categories
            spec["series"] = series
            if spec["kind"] == "pie":
                # 円は項目ごとに色が変わる。値と同じ数だけ色を並べる
                spec["pointColors"] = _point_colors(chart.series[0], len(series[0]["values"]))
            items.append(spec)
        if items:
            by_sheet[name] = items
    return by_sheet


def excel_preview(path: Path) -> dict[str, Any]:
    wb_v = load_workbook(str(path), data_only=True, read_only=False)
    wb_f = load_workbook(str(path), data_only=False, read_only=False)
    images = _excel_images(wb_v, path)
    charts = _excel_charts(wb_v, wb_f)
    sheets = []
    for name in wb_v.sheetnames:
        ws_v, ws_f = wb_v[name], wb_f[name]
        # 画像・グラフはデータの外側(表の右隣など)に置かれることが多い。フロントはアンカーのセルを
        # 基準に重ねるので、そのセルが存在するようにグリッドを広げておく
        overlays = images.get(name, []) + charts.get(name, [])
        anchor_rows = max((o["r"] + 1 for o in overlays), default=1)
        anchor_cols = max((o["c"] + 1 for o in overlays), default=1)
        n_rows = min(max(ws_v.max_row, anchor_rows, 1), MAX_ROWS)
        n_cols = min(max(ws_v.max_column, anchor_cols, 1), MAX_COLS)
        rows = []
        for r in range(1, n_rows + 1):
            row_cells = []
            for c in range(1, n_cols + 1):
                cv = ws_v.cell(row=r, column=c)
                cf = ws_f.cell(row=r, column=c)
                formula = None
                if isinstance(cf.value, str) and cf.value.startswith("="):
                    formula = cf.value
                raw = cv.value if cv.value is not None else (None if formula else cf.value)
                display = _display_value(raw, cv.number_format)
                if display == "" and formula:
                    display = formula  # 未計算の数式は数式文字列を表示
                cell: dict[str, Any] = {}
                if display != "":
                    cell["v"] = display
                if formula:
                    cell["f"] = formula
                style: dict[str, Any] = {}
                font = cv.font
                if font.bold:
                    style["b"] = 1
                if font.italic:
                    style["i"] = 1
                if font.size and float(font.size) != 11.0:
                    style["fs"] = float(font.size)
                fc = _color_hex(font.color) if font.color else None
                if fc and fc != "#000000":
                    style["fc"] = fc
                if cv.fill and cv.fill.patternType == "solid":
                    bg = _color_hex(cv.fill.start_color)
                    if bg:
                        style["bg"] = bg
                if cv.alignment and cv.alignment.horizontal in ("center", "right"):
                    style["ha"] = cv.alignment.horizontal
                elif isinstance(raw, (int, float)) and not isinstance(raw, bool):
                    style["ha"] = "right"
                if style:
                    cell["s"] = style
                row_cells.append(cell)
            rows.append(row_cells)
        merges = []
        for m in ws_v.merged_cells.ranges:
            if m.min_row <= n_rows and m.min_col <= n_cols:
                merges.append({
                    "r": m.min_row - 1,
                    "c": m.min_col - 1,
                    "rs": min(m.max_row, n_rows) - m.min_row + 1,
                    "cs": min(m.max_col, n_cols) - m.min_col + 1,
                })
        col_widths = []
        for c in range(1, n_cols + 1):
            dim = ws_v.column_dimensions.get(get_column_letter(c))
            width = dim.width if dim and dim.width else 8.43
            col_widths.append(round(width * 7.5 + 5))
        sheets.append({
            "name": name,
            "rows": rows,
            "merges": merges,
            "colWidths": col_widths,
            "images": images.get(name, []),
            "charts": charts.get(name, []),
            "truncated": ws_v.max_row > MAX_ROWS or ws_v.max_column > MAX_COLS,
        })
    wb_v.close()
    wb_f.close()
    return {"type": "excel", "sheets": sheets}


def _cache_dir_for(path: Path) -> Path:
    mtime = int(path.stat().st_mtime * 1000)
    return PREVIEW_CACHE_DIR / f"{path.stem}__{mtime}"


def _page_no(p: Path) -> int:
    try:
        return int(p.stem.split("-")[-1])
    except ValueError:
        return 0


def _page_images(path: Path) -> list[Path]:
    """docx/pptxをページPNG群に変換し、ページ順のパスリストを返す(mtimeでキャッシュ)。"""
    cache_dir = _cache_dir_for(path)
    with _soffice_lock:
        if not cache_dir.exists():
            _convert_to_pngs(path, cache_dir)
    images = sorted(cache_dir.glob("slide-*.png"), key=_page_no)
    if not images:
        raise RuntimeError("ページ画像の生成に失敗しました")
    return images


def _pptx_notes(path: Path) -> list[str]:
    """各スライドの発表者ノート本文を、スライド順のリストで返す(ノートが無ければ空文字)。
    (has_notes_slideで確認してから読み、空のノートスライドを生成しないようにする)"""
    prs = Presentation(str(path))
    return [
        slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        for slide in prs.slides
    ]


async def pptx_preview(path: Path) -> dict[str, Any]:
    """PPTXをスライドPNG群に変換し、画像URLと各スライドの発表者ノートを返す。"""
    images = await asyncio.to_thread(_page_images, path)
    notes = await asyncio.to_thread(_pptx_notes, path)
    # 通常はスライド数=画像数だが、変換のずれに備えて画像数に合わせる
    notes = (notes + [""] * len(images))[: len(images)]
    return {
        "type": "pptx",
        "slides": [f"/api/preview_cache/{img.parent.name}/{img.name}" for img in images],
        "notes": notes,
    }


def render_page_png(path: Path, page: int) -> tuple[Path, int]:
    """指定ページ/スライドのPNGパスと総ページ数を返す(エージェントのrender_pageツール用)。
    ページ番号が範囲外のときはValueError。同期関数なのでワーカースレッドから呼べる。"""
    images = _page_images(path)
    if page < 1 or page > len(images):
        raise ValueError(f"ページ番号 {page} は範囲外です (1〜{len(images)})")
    return images[page - 1], len(images)


def _convert_to_pngs(path: Path, cache_dir: Path):
    tmp_dir = cache_dir.with_suffix(".tmp")
    shutil.rmtree(tmp_dir, ignore_errors=True)
    tmp_dir.mkdir(parents=True)
    try:
        subprocess.run(
            [
                "soffice", "--headless", "--norestore",
                f"-env:UserInstallation=file:///tmp/lo_profile",
                "--convert-to", "pdf", "--outdir", str(tmp_dir), str(path),
            ],
            check=True, capture_output=True, timeout=120,
        )
        pdfs = list(tmp_dir.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("PDF変換に失敗しました")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "110", str(pdfs[0]), str(tmp_dir / "slide")],
            check=True, capture_output=True, timeout=120,
        )
        pdfs[0].unlink()
        # 古いキャッシュを掃除してから確定(作業中のtmpディレクトリ自身は除外)。
        # ファイル名に * [ ] ? が含まれてもパターン扱いされないようエスケープする
        for old in PREVIEW_CACHE_DIR.glob(f"{glob_module.escape(path.stem)}__*"):
            if old != tmp_dir:
                shutil.rmtree(old, ignore_errors=True)
        tmp_dir.rename(cache_dir)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
