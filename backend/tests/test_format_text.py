"""キーワードへの文字書式ツール(word/ppt/excel_format_text)のテスト。

「重要な用語を赤字にして」のような依頼で、キーワードの部分だけに色が付き、
前後の文字や既存の書式(太字等)が壊れないことを確認する。
特にWord/PPTでは、キーワードが複数のrunにまたがっている場合の分割処理が肝。
"""
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from docx.oxml.ns import qn as w_qn
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGB
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.oxml.ns import qn as a_qn
from pptx.util import Cm
from pptx.util import Pt as PptPt

from app.agent.tools.excel_tools import excel_format_text
from app.agent.tools.ppt_tools import ppt_format_text
from app.agent.tools.word_tools import word_format_text
from app.atomic import atomic_save

RED = DocxRGB(0xFF, 0x00, 0x00)


# ---------- Word ----------

def _make_docx(ws_dir, build):
    path = ws_dir / "doc.docx"
    doc = Document()
    build(doc)
    atomic_save(doc.save, path)
    return path


def test_word_colors_only_keyword(ws):
    _make_docx(ws, lambda d: d.add_paragraph("機械学習ではファインチューニングが重要です。"))
    result = word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["ファインチューニング"], "color": "#FF0000"}
    )
    assert "1段落・1箇所" in result
    p = Document(str(ws / "doc.docx")).paragraphs[0]
    assert p.text == "機械学習ではファインチューニングが重要です。"  # 本文は変わらない
    red = [r.text for r in p.runs if r.font.color and r.font.color.rgb == RED]
    assert red == ["ファインチューニング"]


def test_word_keyword_split_across_runs(ws):
    def build(d):
        p = d.add_paragraph()
        p.add_run("手法として")
        p.add_run("ファインチ")  # Wordが不規則にrunを分けた状態を再現
        p.add_run("ューニングを使う")

    _make_docx(ws, build)
    result = word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["ファインチューニング"], "color": "#FF0000"}
    )
    assert "1箇所" in result
    p = Document(str(ws / "doc.docx")).paragraphs[0]
    assert p.text == "手法としてファインチューニングを使う"
    red = "".join(r.text for r in p.runs if r.font.color and r.font.color.rgb == RED)
    assert red == "ファインチューニング"
    not_red = "".join(r.text for r in p.runs if not (r.font.color and r.font.color.rgb == RED))
    assert not_red == "手法としてを使う"


def test_word_split_preserves_existing_format(ws):
    def build(d):
        p = d.add_paragraph()
        run = p.add_run("前段RLHF後段")
        run.bold = True

    _make_docx(ws, build)
    word_format_text.invoke({"filename": "doc.docx", "keywords": ["RLHF"], "color": "#FF0000"})
    p = Document(str(ws / "doc.docx")).paragraphs[0]
    assert p.text == "前段RLHF後段"
    assert all(r.bold for r in p.runs)  # 分割で増えたrunも元の太字を引き継ぐ
    red = [r.text for r in p.runs if r.font.color and r.font.color.rgb == RED]
    assert red == ["RLHF"]


def test_word_formats_table_and_multiple_hits(ws):
    def build(d):
        d.add_paragraph("報酬設計と報酬設計")
        table = d.add_table(rows=1, cols=1)
        table.cell(0, 0).text = "報酬設計の一覧"

    _make_docx(ws, build)
    result = word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["報酬設計"], "bold": True}
    )
    assert "2段落" in result and "3箇所" in result
    doc = Document(str(ws / "doc.docx"))
    bold = [r.text for p in doc.paragraphs for r in p.runs if r.bold]
    assert bold == ["報酬設計", "報酬設計"]
    cell_p = doc.tables[0].cell(0, 0).paragraphs[0]
    assert [r.text for r in cell_p.runs if r.bold] == ["報酬設計"]


def test_word_not_found_and_bad_args(ws):
    _make_docx(ws, lambda d: d.add_paragraph("本文"))
    assert "見つかりませんでした" in word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["存在しない語"], "color": "#FF0000"}
    )
    assert "エラー" in word_format_text.invoke({"filename": "doc.docx", "keywords": [], "color": "#FF0000"})
    assert "エラー" in word_format_text.invoke({"filename": "doc.docx", "keywords": ["本文"]})
    assert "エラー" in word_format_text.invoke({"filename": "doc.docx", "keywords": ["本文"], "color": "赤"})
    assert "エラー" in word_format_text.invoke({"filename": "doc.docx", "keywords": ["本文"], "highlight": "黄"})
    assert "エラー" in word_format_text.invoke({"filename": "doc.docx", "keywords": ["本文"], "font_size": 0})


def test_word_font_size_and_font(ws):
    _make_docx(ws, lambda d: d.add_paragraph("要点は結論から書くこと"))
    result = word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["結論"], "font_size": 14, "font": "游ゴシック"}
    )
    assert "1箇所" in result and "14pt" in result and "游ゴシック" in result
    p = Document(str(ws / "doc.docx")).paragraphs[0]
    assert p.text == "要点は結論から書くこと"
    hit = next(r for r in p.runs if r.text == "結論")
    assert hit.font.size == DocxPt(14)
    assert hit.font.name == "游ゴシック"
    # 日本語文字用のeastAsiaフォントも設定されている(これがないと日本語に効かない)
    assert hit._element.rPr.rFonts.get(w_qn("w:eastAsia")) == "游ゴシック"
    other = next(r for r in p.runs if r.text != "結論")
    assert other.font.size is None and other.font.name is None


def test_word_highlight_rounds_to_palette(ws):
    _make_docx(ws, lambda d: d.add_paragraph("ここが重要ポイントです"))
    result = word_format_text.invoke(
        {"filename": "doc.docx", "keywords": ["重要ポイント"], "highlight": "#FFEE55"}  # 黄色に近い色
    )
    assert "1箇所" in result and "#FFFF00" in result and "丸め" in result
    p = Document(str(ws / "doc.docx")).paragraphs[0]
    hit = next(r for r in p.runs if r.text == "重要ポイント")
    assert hit.font.highlight_color == WD_COLOR_INDEX.YELLOW
    assert all(r.font.highlight_color is None for r in p.runs if r.text != "重要ポイント")


# ---------- PowerPoint ----------

def _make_pptx(ws_dir):
    path = ws_dir / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙
    box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(10), Cm(3))
    box.text_frame.text = "強化学習とRLHFの概要"
    table = slide.shapes.add_table(1, 1, Cm(1), Cm(5), Cm(8), Cm(2)).table
    table.cell(0, 0).text = "RLHFの手順"
    atomic_save(prs.save, path)
    return path


def test_ppt_colors_keyword_in_textbox_and_table(ws):
    _make_pptx(ws)
    result = ppt_format_text.invoke(
        {"filename": "deck.pptx", "keywords": ["RLHF"], "color": "#FF0000", "bold": True}
    )
    assert "2箇所" in result
    prs = Presentation(str(ws / "deck.pptx"))
    slide = prs.slides[0]
    box_p = slide.shapes[0].text_frame.paragraphs[0]
    assert "".join(r.text for r in box_p.runs) == "強化学習とRLHFの概要"
    red = [r for r in box_p.runs if r.font.color.type is not None and str(r.font.color.rgb) == "FF0000"]
    assert [r.text for r in red] == ["RLHF"] and red[0].font.bold is True
    cell_p = slide.shapes[1].table.cell(0, 0).text_frame.paragraphs[0]
    assert "".join(r.text for r in cell_p.runs) == "RLHFの手順"
    assert [r.text for r in cell_p.runs if r.font.color.type is not None] == ["RLHF"]


def test_ppt_font_size_font_and_highlight(ws):
    _make_pptx(ws)
    result = ppt_format_text.invoke(
        {"filename": "deck.pptx", "keywords": ["RLHF"],
         "font_size": 28, "font": "メイリオ", "highlight": "#FFFF00"}
    )
    assert "2箇所" in result and "28pt" in result
    prs = Presentation(str(ws / "deck.pptx"))
    box_p = prs.slides[0].shapes[0].text_frame.paragraphs[0]
    assert "".join(r.text for r in box_p.runs) == "強化学習とRLHFの概要"
    hit = next(r for r in box_p.runs if r.text == "RLHF")
    assert hit.font.size == PptPt(28)
    assert hit.font.name == "メイリオ"  # ラテン文字用(a:latin)
    rpr = hit._r.find(a_qn("a:rPr"))
    assert rpr.find(a_qn("a:ea")).get("typeface") == "メイリオ"  # 日本語文字用
    hl = rpr.find(a_qn("a:highlight"))
    assert hl.find(a_qn("a:srgbClr")).get("val") == "FFFF00"
    children = list(rpr)
    assert children.index(hl) < children.index(rpr.find(a_qn("a:latin")))  # スキーマ順が保たれている
    for r in box_p.runs:
        if r.text != "RLHF":
            other_rpr = r._r.find(a_qn("a:rPr"))
            assert other_rpr is None or other_rpr.find(a_qn("a:highlight")) is None


def test_ppt_slide_number_out_of_range(ws):
    _make_pptx(ws)
    result = ppt_format_text.invoke(
        {"filename": "deck.pptx", "keywords": ["RLHF"], "color": "#FF0000", "slide_number": 5}
    )
    assert "範囲外" in result


# ---------- Excel ----------

def test_excel_formats_matching_cells(ws):
    path = ws / "book.xlsx"
    wb = Workbook()
    sh = wb.active
    sh["A1"] = "項目"
    sh["A2"] = "強化学習の説明"
    sh["A3"] = "その他"
    sh["B2"] = '="強化学習"'  # 数式は対象外
    sh["A4"] = "強化学習"
    sh["A4"].font = Font(bold=True)
    atomic_save(wb.save, path)

    result = excel_format_text.invoke(
        {"filename": "book.xlsx", "keywords": ["強化学習"], "color": "#FF0000"}
    )
    assert "2セル" in result and "A2" in result and "A4" in result
    wb2 = load_workbook(str(path))
    sh2 = wb2.active
    assert sh2["A2"].font.color.rgb == "00FF0000"
    assert sh2["A4"].font.color.rgb == "00FF0000"
    assert sh2["A4"].font.bold is True  # 既存の太字は保たれる
    assert sh2["A3"].font.color is None or sh2["A3"].font.color.rgb != "00FF0000"
    assert sh2["B2"].font.color is None or sh2["B2"].font.color.rgb != "00FF0000"


def test_excel_font_size_font_and_highlight(ws):
    path = ws / "book.xlsx"
    wb = Workbook()
    wb.active["A1"] = "強化学習の説明"
    wb.active["A2"] = "その他"
    atomic_save(wb.save, path)

    result = excel_format_text.invoke(
        {"filename": "book.xlsx", "keywords": ["強化学習"],
         "font_size": 14, "font": "游ゴシック", "highlight": "#FFFF00"}
    )
    assert "1セル" in result and "A1" in result
    sh = load_workbook(str(path)).active
    assert sh["A1"].font.size == 14
    assert sh["A1"].font.name == "游ゴシック"
    assert sh["A1"].fill.fill_type == "solid" and sh["A1"].fill.start_color.rgb == "00FFFF00"
    assert sh["A2"].fill.fill_type != "solid"


def test_excel_missing_sheet_and_not_found(ws):
    path = ws / "book.xlsx"
    wb = Workbook()
    wb.active["A1"] = "データ"
    atomic_save(wb.save, path)
    assert "ありません" in excel_format_text.invoke(
        {"filename": "book.xlsx", "keywords": ["データ"], "color": "#FF0000", "sheet": "無いシート"}
    )
    assert "ありませんでした" in excel_format_text.invoke(
        {"filename": "book.xlsx", "keywords": ["無い語"], "color": "#FF0000"}
    )
