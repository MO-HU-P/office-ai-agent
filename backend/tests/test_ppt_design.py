"""おまかせスライドデザイン(同梱テンプレート+デザインガイド注入)のテスト。"""
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from langchain_core.messages import AIMessage, HumanMessage

from app.agent.loop import _load_design_guide, _wants_design_guide
from app.agent.tools.ppt_tools import _TEMPLATE_PATH, _fit_width, ppt_create


def test_bundled_template_exists_and_is_16x9():
    assert _TEMPLATE_PATH.exists()
    prs = Presentation(str(_TEMPLATE_PATH))
    assert len(prs.slides) == 0  # テンプレートにスライドは含まれない
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)


def test_ppt_create_uses_design_template(ws):
    result = ppt_create.invoke({"filename": "deck.pptx", "title": "タイトル", "subtitle": "サブ"})
    assert "作成しました" in result
    prs = Presentation(str(ws / "deck.pptx"))
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)
    theme_xml = prs.slide_master.part.part_related_by(RT.THEME).blob.decode()
    assert '<a:accent1><a:srgbClr val="1A73E8"/></a:accent1>' in theme_xml  # アプリUIと同じ青


def test_fit_width_shrinks_on_narrow_slide(ws):
    # 4:3(幅25.4cm)の既存ファイルでは、16:9向けの既定幅30.5cmが右余白0.5cmまで縮む
    prs = Presentation()  # python-pptx標準は4:3
    assert _fit_width(prs, 1.7, 30.5) == 25.4 - 1.7 - 0.5
    # 収まる場合はそのまま
    ppt_create.invoke({"filename": "wide.pptx"})
    prs_wide = Presentation(str(ws / "wide.pptx"))
    assert _fit_width(prs_wide, 1.7, 30.5) == 30.5


def test_design_guide_loads():
    guide = _load_design_guide()
    assert "スライドデザインガイド" in guide


def test_wants_design_guide_triggers():
    assert _wants_design_guide("営業資料のスライドを作って", [])
    assert _wants_design_guide("presentation.pptx を直して", [])
    assert _wants_design_guide("パワポで3枚にまとめて", [])
    assert not _wants_design_guide("売上をExcelで集計して", [])
    # 直近の履歴にpptx作業があれば「続けて」でも注入される
    history = [HumanMessage("資料.pptx を作って"), AIMessage("作成しました")]
    assert _wants_design_guide("続きをお願い", history)
